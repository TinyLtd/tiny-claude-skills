"""
Native PPTX charts in the Tiny palette.

Every chart here is a real, editable PowerPoint/Google-Slides chart object
(finance can retype the numbers in-place), styled to the house system: blue ->
indigo -> violet -> azure, clean hairline axes, no gridlines, Hanken labels.

Builders return the chart's GraphicFrame. All guard against empty/partial data
so a blank template renders cleanly.
"""
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, XL_TICK_MARK,
    XL_TICK_LABEL_POSITION,
)
from pptx.chart.data import CategoryChartData, BubbleChartData
from pptx.oxml.ns import qn
import copy

import theme as T


def _style_axis(axis, hide_line=False, num_fmt=None, font_pt=11,
                gridlines=False, color=T.MUTED):
    try:
        axis.has_major_gridlines = gridlines
        axis.has_minor_gridlines = False
    except Exception:
        pass
    try:
        axis.major_tick_mark = XL_TICK_MARK.NONE
        axis.minor_tick_mark = XL_TICK_MARK.NONE
    except Exception:
        pass
    tl = axis.tick_labels
    try:
        tl.font.size = Pt(font_pt)
        tl.font.name = T.FONT
        tl.font.color.rgb = color
    except Exception:
        pass
    if num_fmt:
        try:
            tl.number_format = num_fmt
            tl.number_format_is_linked = False
        except Exception:
            pass
    if hide_line:
        try:
            axis.format.line.fill.background()
        except Exception:
            pass
    else:
        try:
            axis.format.line.color.rgb = T.LINE
            axis.format.line.width = Pt(0.75)
        except Exception:
            pass


def _clean(chart, legend=False, legend_pos=XL_LEGEND_POSITION.BOTTOM):
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.name = T.FONT
        chart.legend.font.color.rgb = T.INK


def _color_points(series, colors):
    """Color each data point in a single series individually."""
    for idx, pt in enumerate(series.points):
        c = colors[idx % len(colors)]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = c
        pt.format.line.fill.background()


def _data_labels(plot_or_series, num_fmt="#,##0", pt=11, color=T.INK,
                 pos=XL_LABEL_POSITION.OUTSIDE_END, bold=True):
    try:
        plot_or_series.has_data_labels = True
    except AttributeError:
        pass
    dl = plot_or_series.data_labels
    dl.number_format = num_fmt
    dl.number_format_is_linked = False
    dl.font.size = Pt(pt)
    dl.font.name = T.FONT
    dl.font.bold = bold
    dl.font.color.rgb = color
    try:
        dl.position = pos
    except Exception:
        pass
    return dl


# --------------------------------------------------------------------------
def column_chart(slide, x, y, w, h, cats, values, color=None,
                 per_point_colors=None, num_fmt="#,##0", labels=True,
                 y_fmt="#,##0", show_y=False):
    """Single-series clustered column. per_point_colors colors each bar."""
    cd = CategoryChartData()
    cd.categories = cats or [""]
    cd.add_series("Series", values or [0])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart)
    plot = chart.plots[0]
    plot.gap_width = 55
    plot.vary_by_categories = False
    ser = plot.series[0]
    if per_point_colors:
        _color_points(ser, per_point_colors)
    else:
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = color or T.BLUE
        ser.format.line.fill.background()
    if labels and values:
        _data_labels(plot, num_fmt=num_fmt, pt=10)
    _style_axis(chart.category_axis, hide_line=False, font_pt=10)
    va = chart.value_axis
    va.visible = show_y
    _style_axis(va, hide_line=not show_y, num_fmt=y_fmt, font_pt=10,
                gridlines=False)
    va.has_major_gridlines = False
    return gf


def combo_column_line(slide, x, y, w, h, cats, col_values, line_values,
                      col_color=None, line_color=None,
                      col_fmt="#,##0", line_fmt='0"%"'):
    """Columns (primary axis) + a margin line (secondary axis). Native combo."""
    cd = CategoryChartData()
    cd.categories = cats or [""]
    cd.add_series("Value", col_values or [0])
    cd.add_series("Margin", line_values or [0])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart, legend=False)
    plot = chart.plots[0]
    plot.gap_width = 55
    plot.vary_by_categories = False
    # color the column series
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = col_color or T.BLUE
    plot.series[0].format.line.fill.background()
    _data_labels(plot.series[0], num_fmt=col_fmt, pt=10)

    # configure the PRIMARY axes now, while chart.value_axis is unambiguous
    _style_axis(chart.category_axis, font_pt=10)
    pva = chart.value_axis
    pva.visible = False
    pva.has_major_gridlines = False
    pva.has_minor_gridlines = False

    # --- convert series[1] into a line on a secondary axis via XML ---
    plotArea = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
    barChart = plotArea.find(qn("c:barChart"))
    line_ser = barChart.findall(qn("c:ser"))[1]
    barChart.remove(line_ser)

    sec_cat_id = "555000002"
    sec_val_id = "555000001"

    def _mk(tag, val):
        e = plotArea.makeelement(qn(tag), {}); e.set("val", val); return e

    # build the lineChart group (must sit immediately after barChart, before axes)
    lineChart = plotArea.makeelement(qn("c:lineChart"), {})
    lineChart.append(_mk("c:grouping", "standard"))
    lineChart.append(_mk("c:varyColors", "0"))
    lineChart.append(line_ser)
    lineChart.append(_mk("c:marker", "1"))
    lineChart.append(_mk("c:axId", sec_cat_id))
    lineChart.append(_mk("c:axId", sec_val_id))
    barChart.addnext(lineChart)

    # style the line series: 2.5pt accent line, no curve
    old = line_ser.find(qn("c:spPr"))
    if old is not None:
        line_ser.remove(old)
    spPr = line_ser.makeelement(qn("c:spPr"), {})
    ln = spPr.makeelement(qn("a:ln"), {}); ln.set("w", "31750")  # 2.5pt
    sf = ln.makeelement(qn("a:solidFill"), {})
    clr = sf.makeelement(qn("a:srgbClr"), {})
    lc = line_color or T.color_for(4)
    clr.set("val", "%02X%02X%02X" % (lc[0], lc[1], lc[2]))
    sf.append(clr); ln.append(sf); spPr.append(ln)
    line_ser.find(qn("c:order")).addnext(spPr)
    sm = line_ser.makeelement(qn("c:smooth"), {}); sm.set("val", "0")
    line_ser.append(sm)

    # secondary category axis (hidden) + value axis (right), inserted after the
    # primary axes so all chart groups still precede all axes.
    primary_val = plotArea.findall(qn("c:valAx"))[0]

    secVal = plotArea.makeelement(qn("c:valAx"), {})
    secVal.append(_mk("c:axId", sec_val_id))
    scv = secVal.makeelement(qn("c:scaling"), {})
    scv.append(_mk("c:orientation", "minMax"))
    scv.append(_mk("c:min", "0"))
    secVal.append(scv)
    secVal.append(_mk("c:delete", "0"))
    secVal.append(_mk("c:axPos", "r"))
    nf = secVal.makeelement(qn("c:numFmt"), {})
    nf.set("formatCode", line_fmt); nf.set("sourceLinked", "0"); secVal.append(nf)
    secVal.append(_mk("c:majorTickMark", "none"))
    secVal.append(_mk("c:minorTickMark", "none"))
    secVal.append(_mk("c:tickLblPos", "nextTo"))
    secVal.append(_mk("c:crossAx", sec_cat_id))
    secVal.append(_mk("c:crosses", "max"))
    primary_val.addnext(secVal)

    secCat = plotArea.makeelement(qn("c:catAx"), {})
    secCat.append(_mk("c:axId", sec_cat_id))
    scc = secCat.makeelement(qn("c:scaling"), {})
    scc.append(_mk("c:orientation", "minMax"))
    secCat.append(scc)
    secCat.append(_mk("c:delete", "1"))
    secCat.append(_mk("c:axPos", "b"))
    secCat.append(_mk("c:majorTickMark", "none"))
    secCat.append(_mk("c:minorTickMark", "none"))
    secCat.append(_mk("c:tickLblPos", "none"))
    secCat.append(_mk("c:crossAx", sec_val_id))
    secVal.addnext(secCat)

    return gf


def doughnut(slide, x, y, w, h, labels, values, colors=None, hole=60,
            legend=True):
    cd = CategoryChartData()
    cd.categories = labels or [""]
    cd.add_series("Mix", values or [1])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart, legend=legend, legend_pos=XL_LEGEND_POSITION.RIGHT)
    plot = chart.plots[0]
    try:
        plot.donut_hole_size = hole
    except Exception:
        pass
    ser = plot.series[0]
    cols = colors or T.CHART_PALETTE
    _color_points(ser, cols)
    return gf


def pie(slide, x, y, w, h, labels, values, colors=None, legend=True,
        labels_pct=True):
    cd = CategoryChartData()
    cd.categories = labels or [""]
    cd.add_series("Split", values or [1])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart, legend=legend, legend_pos=XL_LEGEND_POSITION.RIGHT)
    plot = chart.plots[0]
    ser = plot.series[0]
    _color_points(ser, colors or T.CHART_PALETTE)
    if labels_pct:
        dl = plot.data_labels
        dl.show_percentage = True
        dl.number_format = "0%"
        dl.number_format_is_linked = False
        dl.font.size = Pt(11)
        dl.font.name = T.FONT
        dl.font.bold = True
        dl.font.color.rgb = T.WHITE
    return gf


def bar_chart(slide, x, y, w, h, cats, values, color=None,
              per_point_colors=None, num_fmt="#,##0", labels=True):
    """Horizontal bars (BAR_CLUSTERED). Good for debt facilities, proceeds."""
    cd = CategoryChartData()
    cd.categories = cats or [""]
    cd.add_series("Series", values or [0])
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart)
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.vary_by_categories = False
    ser = plot.series[0]
    if per_point_colors:
        _color_points(ser, per_point_colors)
    else:
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = color or T.BLUE
        ser.format.line.fill.background()
    if labels and values:
        _data_labels(plot, num_fmt=num_fmt, pt=10,
                     pos=XL_LABEL_POSITION.OUTSIDE_END)
    _style_axis(chart.category_axis, font_pt=11, hide_line=True)
    va = chart.value_axis
    va.visible = False
    va.has_major_gridlines = False
    _style_axis(va, hide_line=True, font_pt=10)
    return gf


def waterfall(slide, x, y, w, h, labels, deltas, start, end,
              up_color=None, down_color=None):
    """Waterfall via stacked columns: transparent base + visible delta.

    Categories: Start, <each driver>, End. Fully native/editable.
    """
    up_color = up_color or T.BLUE
    down_color = down_color or T.DRAG
    cats = ["Start"] + list(labels) + ["End"]

    base = [0]
    delta = [start]
    colors = [T.color_for(1)]  # start bar (deep indigo)
    running = start
    for v in deltas:
        if v >= 0:
            base.append(running)
            delta.append(v)
            colors.append(up_color)
            running += v
        else:
            running += v
            base.append(running)
            delta.append(-v)
            colors.append(down_color)
    base.append(0)
    delta.append(end)
    colors.append(T.color_for(1))

    cd = CategoryChartData()
    cd.categories = cats
    cd.add_series("Base", base)
    cd.add_series("Delta", delta)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart)
    plot = chart.plots[0]
    plot.gap_width = 40
    # base series: invisible
    base_ser = plot.series[0]
    base_ser.format.fill.background()
    base_ser.format.line.fill.background()
    # delta series: per-point colors
    delta_ser = plot.series[1]
    _color_points(delta_ser, colors)
    _data_labels(delta_ser, num_fmt="#,##0", pt=9,
                 pos=XL_LABEL_POSITION.INSIDE_END, color=T.WHITE)
    _style_axis(chart.category_axis, font_pt=9)
    va = chart.value_axis
    va.visible = False
    va.has_major_gridlines = False
    _style_axis(va, hide_line=True, font_pt=9)
    return gf


def bubble(slide, x, y, w, h, points, x_range=None, y_range=None,
           x_fmt='0"%"', y_fmt='0"%"'):
    """points = [{label,color,x,y,size}]. Native bubble chart."""
    cd = BubbleChartData()
    if not points:
        points = [{"label": "", "x": 0, "y": 0, "size": 1, "color": T.BLUE}]
    for p in points:
        s = cd.add_series(p["label"])
        s.add_data_point(p["x"], p["y"], (p.get("size", 1) or 1) ** 0.5)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BUBBLE, x, y, w, h, cd)
    chart = gf.chart
    _clean(chart, legend=False)
    for i, ser in enumerate(chart.plots[0].series):
        c = points[i].get("color") or T.color_for(i)
        ser.format.fill.solid()
        ser.format.fill.fore_color.rgb = c if isinstance(c, RGBColor) else T.hex_to_rgb(c)
        ser.format.line.fill.background()
    xa = chart.category_axis  # value (x) axis for bubble
    va = chart.value_axis
    if x_range:
        xa.minimum_scale, xa.maximum_scale = x_range
    if y_range:
        va.minimum_scale, va.maximum_scale = y_range
    _style_axis(xa, num_fmt=x_fmt, font_pt=10, gridlines=True)
    _style_axis(va, num_fmt=y_fmt, font_pt=10, gridlines=True)
    xa.has_major_gridlines = True
    va.has_major_gridlines = True
    return gf
