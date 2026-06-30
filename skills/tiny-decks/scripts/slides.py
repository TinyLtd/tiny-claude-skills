"""
Slide layout builders — the Tiny house style as composable layouts.

Two structural systems, mirroring the chosen design directions:
  • B "Electric Hero"  -> cover() / section() : full-bleed blue, white type.
  • C "Split/Editorial" -> a slim blue title rail (kpi/chart slides) or a blue
    header band (wide tables) + white content.

Each builder takes (deck, spec) and adds one slide. spec is a plain dict (see
references/deck-spec.md). All builders tolerate missing/blank fields so the
blank template renders cleanly.
"""
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

import theme as T
import charts as C

RAIL_W = Inches(4.15)
MARGIN = Inches(0.62)
BAND_H = Inches(1.55)


# ---- shared furniture -----------------------------------------------------
def _blank(deck):
    return deck.slides.add_slide(deck.slide_layouts[6])  # blank layout


def _wordmark(slide, x, y, color=T.WHITE, size=20):
    tb, tf = T.textbox(slide, x, y, Inches(2), Inches(0.5))
    T.set_text(tf, "tiny", size, color=color, bold=True, spacing=-0.9)
    return tb


def _eyebrow(tf, text, color=T.MUTED, size=12.5, first=False):
    if text is None:
        return
    if first:
        T.set_text(tf, text.upper(), size, color=color, bold=True,
                   spacing=2.2)
    else:
        T.add_para(tf, text.upper(), size, color=color, bold=True,
                   spacing=2.2, space_before=0)


def _checklist(slide, x, y, w, items, on_blue=False, size=14):
    """The signature blue-checkmark commentary list. items = html-ish strings
    with optional <b>..</b> bold lead-ins."""
    tb, tf = T.textbox(slide, x, y, w, Inches(3.2))
    detail = RGBColor(0xD7, 0xD7, 0xF6) if on_blue else T.MUTED
    lead = T.WHITE if on_blue else T.INK
    check = T.WHITE if on_blue else T.BLUE
    first = True
    for raw in (items or []):
        segs = _parse_segments(raw, lead, detail)
        # prepend a check glyph
        segs = [("✓  ", True, check)] + segs
        if first:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.16
            for seg in segs:
                run = p.add_run()
                run.text = seg[0]
                T._style_run(run, size, seg[2], seg[1])
            first = False
        else:
            p = T.rich_para(tf, segs, size, line=1.16, space_before=9)
    return tb


def _parse_segments(raw, lead, detail):
    """Tiny <b>..</b> parser -> [(text, bold, color)]. Strips other tags."""
    import re
    raw = raw.replace("&amp;", "&").replace("&nbsp;", " ")
    parts = re.split(r"(<b>|</b>)", raw)
    bold = False
    out = []
    for p in parts:
        if p == "<b>":
            bold = True
        elif p == "</b>":
            bold = False
        elif p:
            clean = re.sub(r"<[^>]+>", "", p)
            if clean:
                out.append((clean, bold, lead if bold else detail))
    return out or [(re.sub(r"<[^>]+>", "", raw), False, detail)]


# ---- COVER (Direction B) --------------------------------------------------
def cover(deck, s):
    slide = _blank(deck)
    T.fill_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H, T.BLUE)
    _wordmark(slide, MARGIN, Inches(0.6), size=30)
    tb, tf = T.textbox(slide, MARGIN, Inches(2.35), Inches(11.0), Inches(3.6))
    _eyebrow(tf, s.get("eyebrow", ""), color=RGBColor(0xC9, 0xC9, 0xFA),
             size=14, first=True)
    T.add_para(tf, s.get("title", ""), 60, color=T.WHITE, bold=True,
               spacing=-1.4, line=1.0, space_before=10)
    if s.get("subtitle"):
        T.add_para(tf, s["subtitle"], 21, color=RGBColor(0xDD, 0xDD, 0xFB),
                   line=1.4, space_before=16)
    # footer row
    fb, ff = T.textbox(slide, MARGIN, Inches(6.55), Inches(12.1), Inches(0.5))
    bits = [b for b in [s.get("period"), s.get("entity"), s.get("preparedBy")] if b]
    T.set_text(ff, "   ·   ".join(bits), 14,
               color=RGBColor(0xC9, 0xC9, 0xFA), bold=True)
    if s.get("confidential"):
        cb, cf = T.textbox(slide, Inches(10.3), Inches(6.55), Inches(2.4),
                           Inches(0.5))
        T.set_text(cf, s["confidential"], 13, color=RGBColor(0xC9, 0xC9, 0xFA),
                   bold=True, align=PP_ALIGN.RIGHT, spacing=1.5)
    return slide


# ---- SECTION DIVIDER (Direction B) ---------------------------------------
def section(deck, s):
    slide = _blank(deck)
    T.fill_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H, T.BLUE)
    _wordmark(slide, MARGIN, Inches(0.6), size=24)
    tb, tf = T.textbox(slide, MARGIN, Inches(2.7), Inches(11.5), Inches(2.6),
                       anchor=MSO_ANCHOR.MIDDLE)
    if s.get("number"):
        T.set_text(tf, s["number"], 22, color=RGBColor(0x9A, 0x99, 0xF7),
                   bold=True, spacing=2)
        T.add_para(tf, s.get("title", ""), 54, color=T.WHITE, bold=True,
                   spacing=-1.2, line=1.0, space_before=10)
    else:
        T.set_text(tf, s.get("title", ""), 54, color=T.WHITE, bold=True,
                   spacing=-1.2, line=1.0)
    if s.get("subtitle"):
        T.add_para(tf, s["subtitle"], 19, color=RGBColor(0xDD, 0xDD, 0xFB),
                   line=1.4, space_before=14)
    return slide


# ---- C-rail (left blue title rail) ---------------------------------------
def _rail(slide, eyebrow, title, commentary=None, note=None):
    T.fill_rect(slide, 0, 0, RAIL_W, T.SLIDE_H, T.BLUE)
    _wordmark(slide, MARGIN, Inches(0.55), size=19)
    tb, tf = T.textbox(slide, MARGIN, Inches(1.45), RAIL_W - MARGIN - Inches(0.3),
                       Inches(2.6))
    _eyebrow(tf, eyebrow, color=RGBColor(0xC9, 0xC9, 0xFA), size=12.5, first=True)
    T.add_para(tf, title or "", 33, color=T.WHITE, bold=True, spacing=-0.7,
               line=1.02, space_before=10)
    if commentary:
        _checklist(slide, MARGIN, Inches(4.05), RAIL_W - MARGIN - Inches(0.3),
                   commentary, on_blue=True, size=13)
    if note:
        nb, nf = T.textbox(slide, MARGIN, Inches(6.7), RAIL_W - MARGIN - Inches(0.3),
                           Inches(0.7))
        T.set_text(nf, note, 10.5, color=RGBColor(0xB9, 0xB9, 0xF2), line=1.3)


def _content_x():
    return RAIL_W + Inches(0.55)


def _content_w():
    return T.SLIDE_W - _content_x() - Inches(0.55)


# ---- KPI slide (C-rail + metric callouts + commentary) -------------------
def kpi(deck, s):
    slide = _blank(deck)
    _rail(slide, s.get("eyebrow"), s.get("title"), note=s.get("railNote"))
    x = _content_x()
    w = _content_w()
    cards = s.get("kpis", [])[:4]
    if cards:
        gap = Inches(0.3)
        cw = (w - gap * (len(cards) - 1)) / max(len(cards), 1)
        iw = cw - Inches(0.2) * 2
        vpt = _uniform_value_pt(cards, iw)
        y = Inches(0.85)
        for i, k in enumerate(cards):
            cx = x + (cw + gap) * i
            _kpi_card(slide, cx, y, cw, Inches(2.5), k, val_pt=vpt)
    # commentary below
    if s.get("commentary"):
        _checklist(slide, x, Inches(3.7), w, s["commentary"], size=14)
    if s.get("footnote"):
        fb, ff = T.textbox(slide, x, Inches(6.85), w, Inches(0.5))
        T.set_text(ff, s["footnote"], 10.5, color=T.MUTED, line=1.3)
    return slide


def _value_pt(value, iw, cap=40):
    """Width-aware font size so the big number never wraps inside the card.

    iw = inner text width (EMU). Renderers don't honor no-wrap reliably, so we
    size to actually fit: bold display digits run ~0.74*pt per char.
    """
    n = max(len(str(value)), 1)
    box_pt = (iw / 12700)               # EMU -> points
    fit = box_pt / (0.74 * n)           # largest size that fits on one line
    return max(min(cap, fit), 15)


def _uniform_value_pt(cards, iw, cap=40):
    """One shared value size across a row of cards (smallest that fits all)."""
    if not cards:
        return cap
    return min(_value_pt(c.get("value", "—"), iw, cap) for c in cards)


def _kpi_card(slide, x, y, w, h, k, val_pt=None):
    T.round_rect(slide, x, y, w, h, T.TINT, radius=0.09)
    pad = Inches(0.2)
    iw = w - pad * 2
    # label
    lb, lf = T.textbox(slide, x + pad, y + pad, iw, Inches(0.5))
    T.set_text(lf, k.get("label", "").upper(), 11, color=T.MUTED, bold=True,
               spacing=1.4, line=1.05)
    # big value — own box, sized to fit on one line
    vb, vf = T.textbox(slide, x + pad, y + Inches(0.6), iw, Inches(0.9))
    vf.word_wrap = False
    size = val_pt if val_pt is not None else _value_pt(k.get("value", "—"), iw)
    T.set_text(vf, str(k.get("value", "—")), size,
               color=T.INK, bold=True, spacing=-0.5, line=1.0)
    cur_y = y + Inches(1.5)
    delta = k.get("delta")
    if delta is not None:
        col = T.delta_color(_to_num(delta), k.get("goodUp", True))
        txt = delta if isinstance(delta, str) else T.fmt_pct(delta, sign=True)
        if k.get("deltaLabel"):
            txt = f"{txt}  {k['deltaLabel']}"
        db, df = T.textbox(slide, x + pad, cur_y, iw, Inches(0.4))
        T.set_text(df, txt, 13, color=col, bold=True)
        cur_y = cur_y + Inches(0.36)
    if k.get("sub"):
        sb, sf = T.textbox(slide, x + pad, cur_y, iw, Inches(0.6))
        T.set_text(sf, k["sub"], 11.5, color=T.MUTED, line=1.2)


def _to_num(v):
    if isinstance(v, (int, float)):
        return v
    try:
        return float(str(v).replace("%", "").replace("+", "").replace(",", ""))
    except Exception:
        return None


# ---- TABLE slide (C-band header + native editable table) -----------------
def table(deck, s):
    slide = _blank(deck)
    _band(slide, s.get("eyebrow"), s.get("title"), s.get("bandMeta"))
    cols = s.get("columns", [])
    rows = s.get("rows", [])
    note = s.get("note")
    x = MARGIN
    y = BAND_H + Inches(0.35)
    w = T.SLIDE_W - MARGIN * 2
    avail_h = T.SLIDE_H - y - (Inches(0.75) if note else Inches(0.4))
    _native_table(slide, x, y, w, avail_h, cols, rows,
                  highlight_last_col=s.get("highlightLastCol", False),
                  total_rows=s.get("totalRows", []),
                  variance_cols=s.get("varianceCols", []))
    if note:
        nb, nf = T.textbox(slide, x, T.SLIDE_H - Inches(0.7), w, Inches(0.5))
        T.set_text(nf, note, 11, color=T.MUTED, line=1.3)
    return slide


def _band(slide, eyebrow, title, meta=None):
    T.fill_rect(slide, 0, 0, T.SLIDE_W, BAND_H, T.BLUE)
    _wordmark(slide, MARGIN, Inches(0.32), size=17)
    tb, tf = T.textbox(slide, MARGIN, Inches(0.62), Inches(10.5), Inches(0.85))
    if eyebrow:
        _eyebrow(tf, eyebrow, color=RGBColor(0xC9, 0xC9, 0xFA), size=11.5,
                 first=True)
        T.add_para(tf, title or "", 27, color=T.WHITE, bold=True, spacing=-0.6,
                   space_before=3, line=1.0)
    else:
        T.set_text(tf, title or "", 27, color=T.WHITE, bold=True, spacing=-0.6)
    if meta:
        mb, mf = T.textbox(slide, T.SLIDE_W - Inches(4.0), Inches(0.7),
                           Inches(3.4), Inches(0.6))
        T.set_text(mf, meta, 12, color=RGBColor(0xC9, 0xC9, 0xFA), bold=True,
                   align=PP_ALIGN.RIGHT)


def _native_table(slide, x, y, w, h, cols, rows, highlight_last_col=False,
                  total_rows=None, variance_cols=None):
    total_rows = total_rows or []
    variance_cols = variance_cols or []
    if not cols:
        return
    nrows = len(rows) + 1
    ncols = len(cols)
    gtable = slide.shapes.add_table(nrows, ncols, x, y, w, h).table
    gtable.first_row = False
    gtable.horz_banding = False
    # column widths: first col wide, rest even
    first_w = int(w * 0.22)
    rest = int((w - first_w) / max(ncols - 1, 1))
    gtable.columns[0].width = Emu(first_w)
    for c in range(1, ncols):
        gtable.columns[c].width = Emu(rest)
    # header row
    for c, col in enumerate(cols):
        cell = gtable.cell(0, c)
        _cell(cell, col.get("label", "") if isinstance(col, dict) else str(col),
              size=12, color=T.WHITE, bold=True,
              align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT,
              fill=T.color_for(1))
    # body
    for r, row in enumerate(rows, start=1):
        is_total = (r - 1) in total_rows
        for c in range(ncols):
            val = row[c] if c < len(row) else ""
            cell = gtable.cell(r, c)
            color = T.INK
            bold = is_total or c == 0
            # variance coloring
            if c in variance_cols and isinstance(val, str) and val not in ("", "—"):
                num = _to_num(val)
                if num is not None:
                    color = T.delta_color(num, True)
            fill = T.TINT if is_total else (T.WHITE)
            if highlight_last_col and c == ncols - 1 and not is_total:
                fill = RGBColor(0xEC, 0xEB, 0xFE)
            # brand dot in first column if row carries a color
            _cell(cell, str(val), size=12, color=color, bold=bold,
                  align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT,
                  fill=fill)
    # tighten row heights
    rh = int(h / nrows)
    for r in range(nrows):
        gtable.rows[r].height = Emu(rh)
    return gtable


def _cell(cell, text, size=12, color=T.INK, bold=False, align=PP_ALIGN.LEFT,
          fill=None):
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.03)
    cell.margin_bottom = Inches(0.03)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if fill is None:
        cell.fill.background()
    else:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    T.set_text(tf, text, size, color=color, bold=bold, align=align, line=1.0)


# ---- CHART slide (C-rail + a native chart) -------------------------------
def chart(deck, s):
    slide = _blank(deck)
    _rail(slide, s.get("eyebrow"), s.get("title"), commentary=s.get("commentary"),
          note=s.get("note"))
    x = _content_x()
    w = _content_w()
    y = Inches(0.9)
    h = Inches(5.7)
    _dispatch_chart(slide, x, y, w, h, s)
    return slide


def _dispatch_chart(slide, x, y, w, h, s):
    t = s.get("chartType", "column")
    d = s.get("data", {})
    if t == "column":
        C.column_chart(slide, x, y, w, h, d.get("cats", []), d.get("values", []),
                       color=T.hex_to_rgb(d["color"]) if d.get("color") else None,
                       per_point_colors=_pp(d.get("colors")),
                       num_fmt=d.get("fmt", "#,##0"))
    elif t == "combo":
        C.combo_column_line(slide, x, y, w, h, d.get("cats", []),
                            d.get("values", []), d.get("line", []),
                            col_fmt=d.get("fmt", "#,##0"),
                            line_fmt=d.get("lineFmt", '0"%"'))
    elif t == "doughnut":
        C.doughnut(slide, x, y, w, h, d.get("labels", []), d.get("values", []),
                   colors=_pp(d.get("colors")))
    elif t == "pie":
        C.pie(slide, x, y, w, h, d.get("labels", []), d.get("values", []),
              colors=_pp(d.get("colors")))
    elif t == "bar":
        C.bar_chart(slide, x, y, w, h, d.get("cats", []), d.get("values", []),
                    per_point_colors=_pp(d.get("colors")),
                    num_fmt=d.get("fmt", "#,##0"))
    elif t == "waterfall":
        C.waterfall(slide, x, y, w, h, d.get("labels", []), d.get("deltas", []),
                    d.get("start", 0), d.get("end", 0))
    elif t == "bubble":
        C.bubble(slide, x, y, w, h, d.get("points", []),
                 x_range=d.get("xRange"), y_range=d.get("yRange"),
                 x_fmt=d.get("xFmt", '0"%"'), y_fmt=d.get("yFmt", '0"%"'))


def _pp(colors):
    if not colors:
        return None
    return [T.hex_to_rgb(c) if isinstance(c, str) else c for c in colors]


# ---- DUAL CHART slide (two charts side by side, e.g. twin donuts) ---------
def dualchart(deck, s):
    slide = _blank(deck)
    _band(slide, s.get("eyebrow"), s.get("title"), s.get("bandMeta"))
    half = (T.SLIDE_W - MARGIN * 2 - Inches(0.5)) / 2
    y = BAND_H + Inches(0.45)
    h = Inches(4.0)
    for i, sub in enumerate(s.get("charts", [])[:2]):
        cx = MARGIN + (half + Inches(0.5)) * i
        cb, cf = T.textbox(slide, cx, y - Inches(0.05), half, Inches(0.4))
        T.set_text(cf, sub.get("title", ""), 15, color=T.INK, bold=True,
                   align=PP_ALIGN.CENTER)
        _dispatch_chart(slide, cx, y + Inches(0.45), half, h, sub)
    if s.get("note"):
        nb, nf = T.textbox(slide, MARGIN, T.SLIDE_H - Inches(0.75),
                           T.SLIDE_W - MARGIN * 2, Inches(0.5))
        T.set_text(nf, s["note"], 11, color=T.MUTED, line=1.3, align=PP_ALIGN.CENTER)
    return slide


# ---- STAT GRID (big-number scoreboard) -----------------------------------
def statgrid(deck, s):
    slide = _blank(deck)
    _band(slide, s.get("eyebrow"), s.get("title"), s.get("bandMeta"))
    stats = s.get("stats", [])
    x = MARGIN
    y = BAND_H + Inches(0.55)
    w = T.SLIDE_W - MARGIN * 2
    per_row = s.get("perRow", 3)
    gap = Inches(0.3)
    cw = (w - gap * (per_row - 1)) / per_row
    ch = Inches(1.85)
    iw = cw - Inches(0.2) * 2
    vpt = _uniform_value_pt(stats, iw)
    for i, st in enumerate(stats):
        r, c = divmod(i, per_row)
        cx = x + (cw + gap) * c
        cy = y + (ch + gap) * r
        _kpi_card(slide, cx, cy, cw, ch, st, val_pt=vpt)
    if s.get("commentary"):
        last_row = (len(stats) - 1) // per_row + 1
        cy = y + (ch + gap) * last_row + Inches(0.1)
        _checklist(slide, x, cy, w, s["commentary"], size=13)
    return slide


# ---- GAUGES (leverage / coverage with covenant marker) -------------------
def gauges(deck, s):
    slide = _blank(deck)
    _rail(slide, s.get("eyebrow"), s.get("title"), commentary=s.get("commentary"),
          note=s.get("note"))
    x = _content_x()
    w = _content_w()
    gs = s.get("gauges", [])
    y = Inches(1.0)
    gap = Inches(0.55)
    n = max(len(gs), 1)
    gw = (w - gap * (n - 1)) / n if n > 1 else w * 0.6
    for i, g in enumerate(gs):
        _gauge(slide, x + (gw + gap) * i, y, gw, g)
    return slide


def _gauge(slide, x, y, w, g):
    """A horizontal progress gauge with a covenant threshold marker."""
    label = g.get("label", "")
    val = float(g.get("value", 0) or 0)
    mx = float(g.get("max", 1) or 1)
    marker = g.get("marker")
    unit = g.get("unit", "")
    marker_word = g.get("markerLabel", "covenant")
    frac = max(0.0, min(val / mx, 1.0))
    # big value
    tb, tf = T.textbox(slide, x, y, w, Inches(1.1))
    T.set_text(tf, f"{val:g}{unit}", 46, color=T.BLUE, bold=True, spacing=-1)
    lb, lf = T.textbox(slide, x, y + Inches(1.05), w, Inches(0.5))
    T.set_text(lf, label, 13, color=T.MUTED, bold=True)
    # track
    track_y = y + Inches(1.75)
    track_h = Inches(0.34)
    T.round_rect(slide, x, track_y, w, track_h, T.LINE, radius=0.5)
    if frac > 0:
        T.round_rect(slide, x, track_y, Emu(int(w * frac)), track_h, T.BLUE,
                     radius=0.5)
    # covenant marker
    if marker is not None:
        mfrac = max(0.0, min(float(marker) / mx, 1.0))
        mx_x = x + Emu(int(w * mfrac))
        T.fill_rect(slide, mx_x, track_y - Inches(0.1), Pt(2.2),
                    track_h + Inches(0.2), T.INK)
        mkb, mkf = T.textbox(slide, mx_x - Inches(0.9), track_y + Inches(0.42),
                             Inches(1.8), Inches(0.4))
        T.set_text(mkf, f"{marker_word} {marker:g}{unit}", 10, color=T.MUTED,
                   bold=True, align=PP_ALIGN.CENTER)


# ---- TIMELINE (debt maturities etc.) -------------------------------------
def timeline(deck, s):
    slide = _blank(deck)
    _band(slide, s.get("eyebrow"), s.get("title"), s.get("bandMeta"))
    items = s.get("items", [])
    x = MARGIN
    w = T.SLIDE_W - MARGIN * 2
    y = Inches(4.0)
    T.hline(slide, x, y, w, color=T.LINE, weight=Pt(2))
    n = max(len(items), 1)
    step = w / n
    maxv = max((float(it.get("value", 0) or 0) for it in items), default=1) or 1
    for i, it in enumerate(items):
        cx = x + step * i + step / 2
        v = float(it.get("value", 0) or 0)
        bar_h = Inches(2.4) * (v / maxv)
        bw = Inches(0.9)
        T.round_rect(slide, cx - bw / 2, y - bar_h, bw, bar_h,
                     T.color_for(i), radius=0.12)
        vb, vf = T.textbox(slide, cx - Inches(1), y - bar_h - Inches(0.45),
                           Inches(2), Inches(0.4))
        T.set_text(vf, it.get("valueLabel", T.fmt_money(v)), 13, color=T.INK,
                   bold=True, align=PP_ALIGN.CENTER)
        yb, yf = T.textbox(slide, cx - Inches(1), y + Inches(0.15), Inches(2),
                           Inches(0.4))
        T.set_text(yf, str(it.get("year", "")), 14, color=T.INK, bold=True,
                   align=PP_ALIGN.CENTER)
        lb, lf = T.textbox(slide, cx - Inches(1.1), y + Inches(0.6), Inches(2.2),
                           Inches(0.6))
        T.set_text(lf, it.get("label", ""), 10.5, color=T.MUTED,
                   align=PP_ALIGN.CENTER, line=1.15)
    if s.get("note"):
        nb, nf = T.textbox(slide, x, T.SLIDE_H - Inches(0.7), w, Inches(0.5))
        T.set_text(nf, s["note"], 11, color=T.MUTED, line=1.3)
    return slide


# ---- CONTENT (generic editorial bullets, optional 2-col) -----------------
def content(deck, s):
    slide = _blank(deck)
    _rail(slide, s.get("eyebrow"), s.get("title"), note=s.get("railNote"))
    x = _content_x()
    w = _content_w()
    y = Inches(0.95)
    if s.get("lead"):
        lb, lf = T.textbox(slide, x, y, w, Inches(1.2))
        T.set_text(lf, s["lead"], 20, color=T.INK, bold=True, line=1.2)
        y = Inches(2.2)
    if s.get("bullets"):
        _checklist(slide, x, y, w, s["bullets"], size=15)
    if s.get("footnote"):
        fb, ff = T.textbox(slide, x, T.SLIDE_H - Inches(0.7), w, Inches(0.5))
        T.set_text(ff, s["footnote"], 10.5, color=T.MUTED, line=1.3)
    return slide


# ---- CLOSE (Direction B rally close) -------------------------------------
def close(deck, s):
    slide = _blank(deck)
    T.fill_rect(slide, 0, 0, T.SLIDE_W, T.SLIDE_H, T.BLUE)
    _wordmark(slide, MARGIN, Inches(0.6), size=24)
    tb, tf = T.textbox(slide, MARGIN, Inches(2.6), Inches(11.5), Inches(2.6),
                       anchor=MSO_ANCHOR.MIDDLE)
    T.set_text(tf, s.get("title", ""), 50, color=T.WHITE, bold=True,
               spacing=-1.2, line=1.02)
    if s.get("subtitle"):
        T.add_para(tf, s["subtitle"], 20, color=RGBColor(0xDD, 0xDD, 0xFB),
                   line=1.4, space_before=16)
    if s.get("contact"):
        cb, cf = T.textbox(slide, MARGIN, Inches(6.5), Inches(11), Inches(0.5))
        T.set_text(cf, s["contact"], 14, color=RGBColor(0xC9, 0xC9, 0xFA),
                   bold=True)
    return slide


# ---- confidential watermark ----------------------------------------------
DARK_LAYOUTS = {"cover", "section", "close"}


def add_watermark(slide, layout, text="CONFIDENTIAL"):
    """Faint diagonal watermark across white content slides.

    Skipped on the blue hero layouts (cover/section/close) — they already carry
    the confidential label, and a watermark over a full-bleed hero reads as
    clutter. Light gray so it never fights the data.
    """
    if layout in DARK_LAYOUTS:
        return
    # Small, always-visible corner tag (reliable on dense table slides) ...
    ct, cf = T.textbox(slide, T.SLIDE_W - Inches(2.6), T.SLIDE_H - Inches(0.34),
                       Inches(2.4), Inches(0.3))
    T.set_text(cf, text.upper(), 9, color=RGBColor(0xB6, 0xB6, 0xC8), bold=True,
               align=PP_ALIGN.RIGHT, spacing=2)
    # ... plus a faint diagonal watermark behind the content for the classic look
    tb, tf = T.textbox(slide, Inches(1.2), Inches(2.7), Inches(11.0), Inches(2.0),
                       anchor=MSO_ANCHOR.MIDDLE)
    T.set_text(tf, text.upper(), 60, color=RGBColor(0xEF, 0xEF, 0xF7), bold=True,
               align=PP_ALIGN.CENTER, spacing=3)
    tb.rotation = -22
    sp = tb._element
    tree = sp.getparent()
    tree.remove(sp)
    tree.insert(2, sp)  # behind content


# ---- registry -------------------------------------------------------------
LAYOUTS = {
    "cover": cover,
    "section": section,
    "kpi": kpi,
    "table": table,
    "chart": chart,
    "dualchart": dualchart,
    "statgrid": statgrid,
    "gauges": gauges,
    "timeline": timeline,
    "content": content,
    "close": close,
}
