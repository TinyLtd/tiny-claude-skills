"""
Tiny house style — design tokens and low-level PPTX helpers.

One source of truth for the brand: electric royal-blue (#2B28F5), the heavy
Hanken Grotesk grotesque, the lowercase `tiny` wordmark, giant metric callouts,
blue checkmark commentary, and a cohesive blue->indigo->violet->azure chart
palette (no green/red in charts; green/red is reserved for table variance text).

Everything downstream (slides.py, charts.py) composes from these helpers so the
look stays identical across every deck in the library.
"""
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Canvas ---------------------------------------------------------------
# 16:9 widescreen. python-pptx default is 10x7.5; we use full 13.333x7.5.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ---- Palette --------------------------------------------------------------
BLUE   = RGBColor(0x2B, 0x28, 0xF5)   # signature electric royal-blue
INK    = RGBColor(0x0E, 0x0E, 0x14)   # near-black headlines / body
MUTED  = RGBColor(0x5C, 0x5C, 0x6A)   # gray labels / secondary text
LINE   = RGBColor(0xE5, 0xE5, 0xEE)   # hairline dividers / borders
TINT   = RGBColor(0xF4, 0xF4, 0xFB)   # faint lavender section tint
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PAPER  = RGBColor(0xFF, 0xFF, 0xFF)

# Variance text colors (TABLES ONLY — never in charts)
POS    = RGBColor(0x1A, 0x9E, 0x5C)   # green — favorable delta
NEG    = RGBColor(0xD6, 0x45, 0x45)   # red — unfavorable delta

# Cohesive chart palette: shared chroma, varied hue (blue->indigo->violet->azure)
CHART_PALETTE = [
    RGBColor(0x2B, 0x28, 0xF5),  # electric blue
    RGBColor(0x1A, 0x19, 0x80),  # deep indigo
    RGBColor(0x7A, 0x5C, 0xF5),  # violet
    RGBColor(0x2E, 0x8F, 0xE0),  # azure
    RGBColor(0xC4, 0x4D, 0xDB),  # magenta-violet
    RGBColor(0x5B, 0xB0, 0xE8),  # sky
    RGBColor(0x8E, 0x7B, 0xFF),  # periwinkle
    RGBColor(0xBF, 0xC0, 0xE6),  # pale lavender (drags / NCI / head office)
]
DRAG = RGBColor(0xBF, 0xC0, 0xE6)     # pale-lavender for waterfall declines

# ---- Type -----------------------------------------------------------------
# Hanken Grotesk is a Google Font — renders natively in Google Slides. For
# PowerPoint users without it installed, set FALLBACK; we still name the brand
# font first so Slides looks correct. assets/ ships the TTF for installation.
FONT       = "Hanken Grotesk"
FONT_FALL  = "Arial"

# Projection minimum: nothing below ~14pt at this slide scale (≈24px @1080p).
MIN_PT = 14


def color_for(i):
    """Cycle the chart palette by index."""
    return CHART_PALETTE[i % len(CHART_PALETTE)]


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---- Shape / fill helpers -------------------------------------------------
def _no_line(shape):
    shape.line.fill.background()


def fill_rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    """A solid (or outlined) rectangle. Returns the shape."""
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line_color is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(1)
    return sp


def round_rect(slide, x, y, w, h, color, radius=0.12, line_color=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line_color is None:
        _no_line(sp)
    else:
        sp.line.color.rgb = line_color
        sp.line.width = line_w or Pt(1)
    return sp


def oval(slide, x, y, d, color):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    _no_line(sp)
    return sp


def hline(slide, x, y, w, color=LINE, weight=Pt(0.75)):
    """A hairline divider drawn as a thin connector."""
    conn = slide.shapes.add_connector(2, x, y, x + w, y)  # 2 = straight
    conn.line.color.rgb = color
    conn.line.width = weight
    conn.shadow.inherit = False
    return conn


# ---- Text helpers ---------------------------------------------------------
def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def _style_run(run, size, color, bold, font=FONT, italic=False, spacing=None):
    run.font.size = Pt(max(size, MIN_PT))
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color
    # set the east-asian + complex-script font names too so it sticks
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", font)
    if spacing is not None:
        rPr.set("spc", str(int(spacing * 100)))  # spacing in pt -> 1/100 pt
    return run


def set_text(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, italic=False, spacing=None, line=1.0):
    """Replace a frame's first paragraph with a single styled run."""
    p = tf.paragraphs[0]
    p.alignment = align
    if line:
        p.line_spacing = line
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    run = p.add_run()
    run.text = text
    _style_run(run, size, color, bold, font, italic, spacing)
    return p


def add_para(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT,
             font=FONT, italic=False, spacing=None, line=1.15, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    _style_run(run, size, color, bold, font, italic, spacing)
    return p


def rich_para(tf, segments, size, align=PP_ALIGN.LEFT, line=1.18,
              space_before=0, base_color=MUTED):
    """Add a paragraph from [(text, bold, color), ...] segments.

    Lets a single bullet have a bold lead-in (ink) and gray detail — the
    signature Tiny checkmark-list look.
    """
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line
    if space_before:
        p.space_before = Pt(space_before)
    for seg in segments:
        text, bold, color = seg
        run = p.add_run()
        run.text = text
        _style_run(run, size, color, bold)
    return p


# ---- Number formatting ----------------------------------------------------
def fmt_money(thousands, currency="$", decimals=1, mm=True):
    """Format a $000's figure. mm=True renders as $X.XM; else $X,XXXk."""
    if thousands is None:
        return "—"
    neg = thousands < 0
    v = abs(thousands)
    if mm:
        s = f"{v/1000:.{decimals}f}M"
    else:
        s = f"{v:,.0f}"
    s = f"{currency}{s}"
    return f"({s})" if neg else s


def fmt_pct(v, decimals=0, sign=False):
    if v is None:
        return "—"
    s = f"{v:+.{decimals}f}%" if sign else f"{v:.{decimals}f}%"
    return s


def fmt_num(v, decimals=0):
    if v is None:
        return "—"
    return f"{v:,.{decimals}f}"


def delta_color(v, good_is_up=True):
    if v is None or v == 0:
        return MUTED
    up = v > 0
    favorable = up if good_is_up else (not up)
    return POS if favorable else NEG
